from pathlib import Path

from pptx import Presentation

from app.parsers.base import ParsedBlock


class PptxParser:
    name = "python-pptx"
    version = "1"

    def parse(self, path: Path) -> list[ParsedBlock]:
        presentation = Presentation(path)
        blocks: list[ParsedBlock] = []
        for slide_number, slide in enumerate(presentation.slides, start=1):
            texts: list[str] = []
            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False):
                    value = shape.text.strip()
                    if value:
                        texts.append(value)
                if getattr(shape, "has_table", False):
                    for row in shape.table.rows:
                        value = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                        if value:
                            texts.append(value)
            try:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    texts.append(f"备注：{notes}")
            except (AttributeError, KeyError):
                pass
            if texts:
                title = slide.shapes.title.text.strip() if slide.shapes.title and slide.shapes.title.text else f"幻灯片 {slide_number}"
                blocks.append(ParsedBlock(content="\n".join(texts), slide_number=slide_number, section_path=[title]))
        return blocks

